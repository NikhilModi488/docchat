"""
app/ingestion/loader.py
=======================
Document loading layer — multi-format.

Turns a file on disk into a list of page-level LangChain `Document` objects with
`source` (filename) and `page_number` (1-based) metadata.

Supported formats (spec): PDF, DOCX, TXT, Markdown, PPTX.

PDF gets the richest treatment:
    * PyMuPDF text extraction (per page).
    * Optional OCR fallback (pytesseract) for image-only / scanned pages, used
      only when a page yields almost no text AND Tesseract is available.
    * Optional table extraction (PyMuPDF find_tables) appended as TSV text.

Design (Open/Closed): adding a format = adding a `_load_<ext>` method; the
public `load()` dispatch and the rest of the pipeline stay unchanged.
"""

from __future__ import annotations

from pathlib import Path

from langchain_core.documents import Document

from app.config import settings
from app.utils.logger import get_logger

logger = get_logger(__name__)

# --- optional heavy deps, imported lazily/guarded -------------------------
try:
    import fitz  # PyMuPDF
    _HAS_FITZ = True
except Exception:  # pragma: no cover
    try:
        import pymupdf as fitz  # type: ignore
        _HAS_FITZ = True
    except Exception:
        _HAS_FITZ = False


def _ocr_available() -> bool:
    try:
        import pytesseract  # noqa: F401
        from PIL import Image  # noqa: F401

        # Verify the tesseract binary is actually reachable.
        import pytesseract as pt

        pt.get_tesseract_version()
        return True
    except Exception:
        return False


class DocumentLoader:
    """Loads a file of any supported type into page-level Documents."""

    def load(self, path: Path) -> list[Document]:
        path = Path(path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {path}")

        ext = path.suffix.lower()
        logger.info("Loading %s (%s)", path.name, ext)

        if ext == ".pdf":
            docs = self._load_pdf(path)
        elif ext == ".docx":
            docs = self._load_docx(path)
        elif ext == ".pptx":
            docs = self._load_pptx(path)
        elif ext in (".txt", ".md", ".markdown"):
            docs = self._load_text(path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

        for d in docs:
            d.metadata.setdefault("source", path.name)
        logger.info("Loaded %d page(s) from %s", len(docs), path.name)
        return docs

    # ------------------------------------------------------------------ #
    # PDF
    # ------------------------------------------------------------------ #
    def _load_pdf(self, path: Path) -> list[Document]:
        if _HAS_FITZ:
            return self._load_pdf_fitz(path)
        # Fallback to PyPDF if PyMuPDF missing.
        from langchain_community.document_loaders import PyPDFLoader

        pages = PyPDFLoader(str(path)).load()
        for p in pages:
            p.metadata["source"] = path.name
            p.metadata["page_number"] = int(p.metadata.get("page", 0)) + 1
        return pages

    def _load_pdf_fitz(self, path: Path) -> list[Document]:
        cfg = settings.ingestion
        ocr_ok = cfg.enable_ocr and _ocr_available()
        docs: list[Document] = []
        doc = fitz.open(str(path))
        try:
            for i in range(doc.page_count):
                page = doc.load_page(i)
                text = page.get_text("text") or ""

                # OCR fallback for near-empty (scanned) pages.
                if ocr_ok and len(text.strip()) < 20:
                    try:
                        text = self._ocr_page(page) or text
                    except Exception:
                        logger.exception("OCR failed on page %d of %s", i + 1, path.name)

                # Table extraction appended as TSV.
                if cfg.enable_tables:
                    try:
                        text += self._extract_tables(page)
                    except Exception:
                        logger.debug("No tables / table extraction failed on page %d", i + 1)

                docs.append(
                    Document(
                        page_content=text,
                        metadata={"source": path.name, "page_number": i + 1},
                    )
                )
        finally:
            doc.close()
        return docs

    @staticmethod
    def _ocr_page(page) -> str:
        import io

        import pytesseract
        from PIL import Image

        pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
        img = Image.open(io.BytesIO(pix.tobytes("png")))
        return pytesseract.image_to_string(img)

    @staticmethod
    def _extract_tables(page) -> str:
        finder = page.find_tables()
        out = ""
        for ti, table in enumerate(getattr(finder, "tables", []) or []):
            rows = table.extract()
            if not rows:
                continue
            out += f"\n\n[Table {ti + 1}]\n"
            for row in rows:
                out += "\t".join("" if c is None else str(c) for c in row) + "\n"
        return out

    # ------------------------------------------------------------------ #
    # DOCX
    # ------------------------------------------------------------------ #
    def _load_docx(self, path: Path) -> list[Document]:
        import docx  # python-docx

        d = docx.Document(str(path))
        parts = [p.text for p in d.paragraphs if p.text.strip()]
        # Append tables as TSV.
        for table in d.tables:
            for row in table.rows:
                parts.append("\t".join(cell.text for cell in row.cells))
        text = "\n".join(parts)
        return [Document(page_content=text, metadata={"source": path.name, "page_number": 1})]

    # ------------------------------------------------------------------ #
    # PPTX
    # ------------------------------------------------------------------ #
    def _load_pptx(self, path: Path) -> list[Document]:
        from pptx import Presentation

        prs = Presentation(str(path))
        docs: list[Document] = []
        for idx, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
            content = "\n".join(t for t in texts if t.strip())
            if content.strip():
                docs.append(
                    Document(page_content=content, metadata={"source": path.name, "page_number": idx})
                )
        return docs or [Document(page_content="", metadata={"source": path.name, "page_number": 1})]

    # ------------------------------------------------------------------ #
    # TXT / Markdown
    # ------------------------------------------------------------------ #
    def _load_text(self, path: Path) -> list[Document]:
        text = path.read_text(encoding="utf-8", errors="ignore")
        return [Document(page_content=text, metadata={"source": path.name, "page_number": 1})]
